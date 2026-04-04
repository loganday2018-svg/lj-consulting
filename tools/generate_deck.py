"""
Day Horrigan — Pitch Deck Generator

Usage:
    python3 generate_deck.py --firm "Firm Name" --config firms/firm_name.json

Config JSON format:
{
    "firm_name": "Natural Capital",
    "aum": "$218M",
    "total_investments": 28,
    "total_invested": "$326M",
    "logo_path": "assets/nc_logo.png",
    "portcos": [
        {
            "name": "Dispatch",
            "subtitle": "AI-powered last-mile delivery • Minneapolis",
            "question": "How much of driver-to-route matching still involves manual decisions?",
            "accent_color": [91, 141, 217],
            "bg_shade": "dark",
            "logo_path": "assets/dispatch_webclip.png",
            "logo_w": 0.65,
            "logo_h": 0.65
        }
    ]
}

Colors: Each portco gets its own accent. Use bg_shade "dark" (0x3A3A3A) or "light" (0x444444) to alternate.
Pick 3 portcos max for the deck.
"""

import json
import argparse
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Brand constants ──
CREAM = RGBColor(0xF5, 0xF0, 0xE8)
DARK = RGBColor(0x2C, 0x2C, 0x2C)
ACCENT = RGBColor(0x7E, 0xB0, 0x7E)
LIGHT_ACCENT = RGBColor(0xA8, 0xCC, 0xA8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x6B, 0x6B, 0x6B)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
ASSETS_DIR = os.path.join(SCRIPT_DIR, "assets")
DOT_GRID = os.path.join(ASSETS_DIR, "dot_grid.png")


def set_slide_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def add_bg_texture(slide):
    if os.path.exists(DOT_GRID):
        slide.shapes.add_picture(DOT_GRID, Inches(0), Inches(0), Inches(13.333), Inches(7.5))


def add_shape(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18, color=CREAM, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf


def add_para(tf, text, font_size=18, color=CREAM, bold=False, space_before=Pt(8), alignment=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_before = space_before
    p.alignment = alignment
    return p


def generate_deck(config_path, output_path=None):
    global ACCENT, LIGHT_ACCENT

    with open(config_path) as f:
        config = json.load(f)

    firm = config["firm_name"]
    portcos = config["portcos"][:3]  # Max 3

    # Override accent colors if firm specifies them
    if "accent_color" in config:
        ACCENT = RGBColor(*config["accent_color"])
    if "light_accent_color" in config:
        LIGHT_ACCENT = RGBColor(*config["light_accent_color"])

    if output_path is None:
        safe_name = firm.replace(" ", "_")
        output_path = f"Day_Horrigan_x_{safe_name}.pptx"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # ── SLIDE 1: TITLE ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK)
    add_bg_texture(slide)
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(2.8), CREAM)

    logo_path = config.get("logo_path")
    if logo_path and os.path.exists(logo_path):
        add_shape(slide, Inches(9.2), Inches(0.3), Inches(3.8), Inches(2.0), DARK)
        slide.shapes.add_picture(logo_path, Inches(9.5), Inches(0.55), Inches(3.2), Inches(1.39))

    add_text_box(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(1.2),
                 "Day Horrigan", 44, DARK, True)
    add_text_box(slide, Inches(0.8), Inches(1.6), Inches(8), Inches(0.8),
                 "Turn AI into EBITDA.", 24, ACCENT)
    add_shape(slide, Inches(0.8), Inches(3.5), Inches(1.5), Inches(0.06), ACCENT)
    add_text_box(slide, Inches(0.8), Inches(3.9), Inches(11), Inches(0.6),
                 f"Prepared for {firm}", 22, CREAM)
    add_text_box(slide, Inches(0.8), Inches(4.6), Inches(11), Inches(0.6),
                 "AI Implementation & Automation for the Heartland", 16, MUTED)
    add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.6),
                 "dayhorrigan.vercel.app", 14, ACCENT)

    # ── SLIDE 2: PITCH ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK)
    add_bg_texture(slide)
    tf = add_text_box(slide, Inches(0.8), Inches(1.5), Inches(10), Inches(3.0),
                 "Companies are struggling to realize the full value of AI.", 30, CREAM)
    add_para(tf, "", 14, CREAM)
    add_para(tf, "If your portcos are only using chat features, they're leaving most of the opportunity on the table.", 30, CREAM)
    add_para(tf, "", 14, CREAM)
    add_para(tf, "We fix that. Weeks, not quarters.", 26, ACCENT, True)

    # ── SLIDE 3: BIG STAT ──
    total_inv = config.get("total_investments", "")
    total_amt = config.get("total_invested", "")
    if total_inv or total_amt:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_bg(slide, DARK)
        add_bg_texture(slide)
        if total_inv:
            add_text_box(slide, Inches(1.0), Inches(1.2), Inches(11), Inches(1.2),
                         f"{total_inv} portfolio companies.", 52, WHITE, True, PP_ALIGN.CENTER)
        if total_amt:
            add_text_box(slide, Inches(1.0), Inches(2.6), Inches(11), Inches(1.2),
                         f"{total_amt} invested.", 52, WHITE, True, PP_ALIGN.CENTER)
        add_shape(slide, Inches(6.0), Inches(4.2), Inches(1.3), Inches(0.05), ACCENT)
        add_text_box(slide, Inches(2.0), Inches(4.8), Inches(9.3), Inches(1.0),
                     "How many are using AI beyond chat?", 28, MUTED, False, PP_ALIGN.CENTER)

    # ── SLIDE 4: DIVIDER ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK)
    add_bg_texture(slide)
    add_shape(slide, Inches(6.0), Inches(2.8), Inches(1.3), Inches(0.05), ACCENT)
    add_text_box(slide, Inches(1.0), Inches(3.2), Inches(11.3), Inches(1.0),
                 "Your portfolio. Some questions.", 40, CREAM, True, PP_ALIGN.CENTER)

    # ── SLIDE 5: PORTFOLIO ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK)
    add_bg_texture(slide)

    tops = [Inches(0.6), Inches(2.7), Inches(4.8)]
    for i, pc in enumerate(portcos):
        top = tops[i]
        accent = RGBColor(*pc["accent_color"])
        bg = RGBColor(0x44, 0x44, 0x44) if pc.get("bg_shade") == "light" else RGBColor(0x3A, 0x3A, 0x3A)

        add_shape(slide, Inches(0.5), top, Inches(12.3), Inches(1.9), bg)
        add_shape(slide, Inches(0.5), top, Inches(0.08), Inches(1.9), accent)

        pc_logo = pc.get("logo_path")
        if pc_logo and os.path.exists(pc_logo):
            slide.shapes.add_picture(
                pc_logo,
                Inches(0.9), top + Inches(0.2),
                Inches(pc.get("logo_w", 1.5)), Inches(pc.get("logo_h", 0.5))
            )
            add_text_box(slide, Inches(0.9), top + Inches(0.75), Inches(3.5), Inches(0.4),
                         pc["subtitle"], 12, accent)
        else:
            add_text_box(slide, Inches(0.9), top + Inches(0.15), Inches(3.5), Inches(0.5),
                         pc["name"], 20, WHITE, True)
            add_text_box(slide, Inches(0.9), top + Inches(0.6), Inches(3.5), Inches(0.4),
                         pc["subtitle"], 12, accent)

        add_text_box(slide, Inches(5.5), top + Inches(0.15), Inches(6.8), Inches(1.6),
                     pc["question"], 14, LIGHT_GRAY)

    # ── SLIDE 6: DIVIDER ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK)
    add_bg_texture(slide)
    add_shape(slide, Inches(6.0), Inches(2.8), Inches(1.3), Inches(0.05), ACCENT)
    add_text_box(slide, Inches(1.0), Inches(3.2), Inches(11.3), Inches(1.0),
                 "What we've seen.", 40, CREAM, True, PP_ALIGN.CENTER)

    # ── SLIDE 7: RANDY'S TIMELINE ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK)
    add_bg_texture(slide)
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.4), CREAM)
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
                 "Randy's Worldwide", 32, DARK, True)
    add_text_box(slide, Inches(0.8), Inches(0.85), Inches(11), Inches(0.4),
                 "Consumer products • Multi-department rollout", 14, ACCENT)

    timeline_y = Inches(2.2)
    add_shape(slide, Inches(1.0), timeline_y + Inches(0.35), Inches(11.3), Inches(0.06), ACCENT)
    milestones = [
        ("Week 1", "COO trained\non Claude", Inches(1.5)),
        ("Week 3", "KPI dashboard\nlive from raw data", Inches(4.0)),
        ("Week 6", "P&L analysis\nautomated", Inches(6.5)),
        ("Week 8", "CFO, VP Procurement\nadopting independently", Inches(9.0)),
        ("Week 12", "HR, Product Dev\nusing AI organically", Inches(11.2)),
    ]
    for label, desc, left in milestones:
        add_shape(slide, left - Inches(0.04), timeline_y + Inches(0.11), Inches(0.24), Inches(0.24), ACCENT)
        add_text_box(slide, left - Inches(0.5), timeline_y - Inches(0.4), Inches(1.2), Inches(0.4),
                     label, 13, ACCENT, True, PP_ALIGN.CENTER)
        add_text_box(slide, left - Inches(0.6), timeline_y + Inches(0.6), Inches(1.5), Inches(0.8),
                     desc, 12, LIGHT_GRAY, False, PP_ALIGN.CENTER)

    add_shape(slide, Inches(0.8), Inches(4.5), Inches(11.7), Inches(2.5), RGBColor(0x3A, 0x3A, 0x3A))
    tf = add_text_box(slide, Inches(1.3), Inches(4.7), Inches(5.0), Inches(2.0),
                 "What we did:", 18, WHITE, True)
    add_para(tf, "", 6, WHITE)
    for item in ["• KPI dashboard from raw financials", "• AI deployed across C-suite — COO, CFO, VP Procurement",
                 "• Automated P&L analysis with exec summaries", "• Adoption spread organically to HR, Product Dev, Finance"]:
        add_para(tf, item, 14, LIGHT_GRAY)
    tf2 = add_text_box(slide, Inches(7.0), Inches(4.7), Inches(5.0), Inches(2.0),
                 "< 3 months", 32, WHITE, True)
    add_para(tf2, "Skeptical to self-sufficient", 14, MUTED)

    # ── SLIDE 8: PULL QUOTE ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK)
    add_bg_texture(slide)
    add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10.3), Inches(0.8),
                 "\u201C", 120, ACCENT)
    add_text_box(slide, Inches(2.0), Inches(2.5), Inches(9.3), Inches(2.0),
                 "I am really excited about what\nI just whipped together.", 44, WHITE, True, PP_ALIGN.CENTER)
    add_shape(slide, Inches(6.0), Inches(5.0), Inches(1.3), Inches(0.05), ACCENT)
    add_text_box(slide, Inches(2.0), Inches(5.4), Inches(9.3), Inches(0.6),
                 "COO, Randy's Worldwide", 18, MUTED, False, PP_ALIGN.CENTER)
    add_text_box(slide, Inches(2.0), Inches(5.9), Inches(9.3), Inches(0.5),
                 "After 8 weeks of working with Day Horrigan", 14, MUTED, False, PP_ALIGN.CENTER)

    # ── SLIDE 9: WHO WE ARE ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, DARK)
    add_bg_texture(slide)
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.4), CREAM)
    add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
                 "Who we are", 32, DARK, True)

    for name, sub, bio, left in [
        ("Logan Day", "Darden MBA • U.S. Army Captain",
         "Finance at Walmart by day, where I build AI tools and automation for executive teams. Built the dashboards and AI workflows at Randy's Worldwide.\n\nThis is a focused side practice — we do this because we're genuinely good at it and the demand found us.",
         Inches(0.8)),
        ("Matt Horrigan", "Darden MBA • Strategy & Finance",
         "Strategy and finance at WEX, with corporate strategy experience at Comcast.\n\nBrings the strategic lens — he sees the portco P&L the way your partners do.",
         Inches(7.0)),
    ]:
        add_shape(slide, left, Inches(2.0), Inches(5.5), Inches(4.5), RGBColor(0x3A, 0x3A, 0x3A))
        add_shape(slide, left, Inches(2.0), Inches(0.08), Inches(4.5), ACCENT)
        tf = add_text_box(slide, left + Inches(0.5), Inches(2.3), Inches(4.5), Inches(4.0),
                     name, 24, WHITE, True)
        add_para(tf, sub, 14, LIGHT_ACCENT)
        add_para(tf, "", 10, WHITE)
        add_para(tf, bio, 15, LIGHT_GRAY)

    # ── SLIDE 10: CTA ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, CREAM)
    if logo_path and os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(10.5), Inches(6.0), Inches(2.0), Inches(0.87))
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.2),
                 "30 minutes.", 52, DARK, True, PP_ALIGN.CENTER)
    add_shape(slide, Inches(6.0), Inches(3.0), Inches(1.3), Inches(0.05), ACCENT)
    add_text_box(slide, Inches(2.0), Inches(3.5), Inches(9.3), Inches(0.6),
                 "That's all it takes to find the quick wins.", 20, MUTED, False, PP_ALIGN.CENTER)
    tf = add_text_box(slide, Inches(2.0), Inches(4.5), Inches(9.3), Inches(2.5),
                 "logan.day2018@gmail.com", 22, DARK, True, PP_ALIGN.CENTER)
    add_para(tf, "", 14, DARK)
    add_para(tf, "dayhorrigan.vercel.app", 16, ACCENT, False, Pt(12), PP_ALIGN.CENTER)
    add_para(tf, "", 10, DARK)
    add_para(tf, "Logan Day  •  Matt Horrigan", 16, MUTED, False, Pt(12), PP_ALIGN.CENTER)

    prs.save(output_path)
    print(f"Saved to {output_path}")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generate Day Horrigan pitch deck for a PE firm")
    parser.add_argument("--config", required=True, help="Path to firm config JSON")
    parser.add_argument("--output", help="Output .pptx path")
    args = parser.parse_args()
    generate_deck(args.config, args.output)
