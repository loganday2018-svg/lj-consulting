from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Inverted: dark backgrounds, cream/light text
CREAM = RGBColor(0xF5, 0xF0, 0xE8)
DARK = RGBColor(0x2C, 0x2C, 0x2C)
ACCENT = RGBColor(0x7E, 0xB0, 0x7E)
LIGHT_ACCENT = RGBColor(0xA8, 0xCC, 0xA8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x6B, 0x6B, 0x6B)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)

DISPATCH_ACCENT = RGBColor(0x5B, 0x8D, 0xD9)
VIP_ACCENT = RGBColor(0x4E, 0xB0, 0xA0)
HARVEST_ACCENT = RGBColor(0xD4, 0xA0, 0x4E)

ASSETS = "/Users/loganday/Desktop/deck_assets"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_bg_texture(slide):
    """Add dot grid texture as background image"""
    slide.shapes.add_picture(
        f"{ASSETS}/dot_grid.png",
        Inches(0), Inches(0),
        Inches(13.333), Inches(7.5)
    )

def add_shape(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=CREAM, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf

def add_para(tf, text, font_size=18, color=CREAM, bold=False, space_before=Pt(8), alignment=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_before = space_before
    p.alignment = alignment
    return p

# ═══════════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK)
add_bg_texture(slide)

# Cream header bar
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(2.8), CREAM)

# Natural Capital logo on dark rounded pill (top right)
add_shape(slide, Inches(9.2), Inches(0.3), Inches(3.8), Inches(2.0), DARK)
slide.shapes.add_picture(
    f"{ASSETS}/nc_logo.png",
    Inches(9.5), Inches(0.55),
    Inches(3.2), Inches(1.39)
)

add_text_box(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(1.2),
             "Day Horrigan", 44, DARK, True, PP_ALIGN.LEFT)
add_text_box(slide, Inches(0.8), Inches(1.6), Inches(8), Inches(0.8),
             "Turn AI into EBITDA.", 24, ACCENT, False, PP_ALIGN.LEFT)

# Accent bar
add_shape(slide, Inches(0.8), Inches(3.5), Inches(1.5), Inches(0.06), ACCENT)

add_text_box(slide, Inches(0.8), Inches(3.9), Inches(11), Inches(0.6),
             "Prepared for Natural Capital", 22, CREAM, False)
add_text_box(slide, Inches(0.8), Inches(4.6), Inches(11), Inches(0.6),
             "AI Implementation & Automation for the Heartland", 16, MUTED, False)

add_text_box(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.6),
             "lj-consulting.vercel.app", 14, ACCENT, False)

# ═══════════════════════════════════════════
# SLIDE 2: THE PITCH
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK)
add_bg_texture(slide)

tf = add_text_box(slide, Inches(0.8), Inches(1.5), Inches(10), Inches(3.0),
             "Companies are struggling to realize the full value of AI.", 30, CREAM, False)
add_para(tf, "", 14, CREAM)
add_para(tf, "If your portcos are only using chat features, they're leaving most of the opportunity on the table.", 30, CREAM)
add_para(tf, "", 14, CREAM)
add_para(tf, "We fix that. Weeks, not quarters.", 26, ACCENT, True)

# ═══════════════════════════════════════════
# SLIDE 3: BIG STAT (new)
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK)
add_bg_texture(slide)

add_text_box(slide, Inches(1.0), Inches(1.2), Inches(11), Inches(1.2),
             "28 portfolio companies.", 52, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.0), Inches(2.6), Inches(11), Inches(1.2),
             "$326M invested.", 52, WHITE, True, PP_ALIGN.CENTER)

add_shape(slide, Inches(6.0), Inches(4.2), Inches(1.3), Inches(0.05), ACCENT)

add_text_box(slide, Inches(2.0), Inches(4.8), Inches(9.3), Inches(1.0),
             "How many are using AI beyond chat?", 28, MUTED, False, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# SLIDE 4: DIVIDER — "Your portfolio"
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK)
add_bg_texture(slide)

add_shape(slide, Inches(6.0), Inches(2.8), Inches(1.3), Inches(0.05), ACCENT)
add_text_box(slide, Inches(1.0), Inches(3.2), Inches(11.3), Inches(1.0),
             "Your portfolio. Some questions.", 40, CREAM, True, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# SLIDE 5: PORTFOLIO COMPANIES
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK)
add_bg_texture(slide)

portcos = [
    {
        "top": Inches(0.6),
        "bg": RGBColor(0x3A, 0x3A, 0x3A),
        "accent": DISPATCH_ACCENT,
        "name": "Dispatch",
        "sub": "Last-mile delivery • Minneapolis",
        "question": "How much of driver-to-route matching still involves manual decisions? Are customer delivery updates and exception handling eating ops bandwidth? Is anyone building internal dashboards, or is the team still living in spreadsheets?",
        "logo": f"{ASSETS}/dispatch_webclip.png",
        "logo_w": Inches(0.65),
        "logo_h": Inches(0.65),
    },
    {
        "top": Inches(2.7),
        "bg": RGBColor(0x44, 0x44, 0x44),
        "accent": VIP_ACCENT,
        "name": "Vision Integrated Partners",
        "sub": "23 practices • 56 locations • 12 surgery centers",
        "question": "How many hours per week do staff spend on insurance verification and prior auth across 56 locations? Is cross-location financial reporting still a manual consolidation exercise? Could clinical note summarization free up physician time?",
        "logo": f"{ASSETS}/vip_logo.png",
        "logo_w": Inches(2.2),
        "logo_h": Inches(0.47),
    },
    {
        "top": Inches(4.8),
        "bg": RGBColor(0x3A, 0x3A, 0x3A),
        "accent": HARVEST_ACCENT,
        "name": "The Harvest Group",
        "sub": "Retail marketing • Rogers, AR",
        "question": "How long does it take to pull, format, and deliver campaign reports to clients? Could performance data generate the first draft of creative briefs? What if client-facing dashboards took hours instead of weeks?",
        "logo": f"{ASSETS}/harvest_logo.png",
        "logo_w": Inches(1.6),
        "logo_h": Inches(0.78),
    },
]

for p in portcos:
    # Card background
    add_shape(slide, Inches(0.5), p["top"], Inches(12.3), Inches(1.9), p["bg"])
    # Left accent bar
    add_shape(slide, Inches(0.5), p["top"], Inches(0.08), Inches(1.9), p["accent"])
    # Company logo
    slide.shapes.add_picture(
        p["logo"],
        Inches(0.9), p["top"] + Inches(0.2),
        p["logo_w"], p["logo_h"]
    )
    # Subtitle
    add_text_box(slide, Inches(0.9), p["top"] + Inches(0.75), Inches(3.5), Inches(0.4),
                 p["sub"], 12, p["accent"])
    # Question
    add_text_box(slide, Inches(5.5), p["top"] + Inches(0.15), Inches(6.8), Inches(1.6),
                 p["question"], 14, LIGHT_GRAY)

# ═══════════════════════════════════════════
# SLIDE 6: DIVIDER — "What we've seen."
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK)
add_bg_texture(slide)

add_shape(slide, Inches(6.0), Inches(2.8), Inches(1.3), Inches(0.05), ACCENT)
add_text_box(slide, Inches(1.0), Inches(3.2), Inches(11.3), Inches(1.0),
             "What we've seen.", 40, CREAM, True, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# SLIDE 7: CASE STUDY — TIMELINE
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK)
add_bg_texture(slide)

add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.4), CREAM)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
             "Randy's Worldwide", 32, DARK, True)
add_text_box(slide, Inches(0.8), Inches(0.85), Inches(11), Inches(0.4),
             "Consumer products • Multi-department rollout", 14, ACCENT)

# Timeline bar
timeline_y = Inches(2.2)
add_shape(slide, Inches(1.0), timeline_y + Inches(0.35), Inches(11.3), Inches(0.06), ACCENT)

# Timeline nodes
milestones = [
    ("Week 1", "COO trained\non Claude", Inches(1.5)),
    ("Week 3", "KPI dashboard\nlive from raw data", Inches(4.0)),
    ("Week 6", "P&L analysis\nautomated", Inches(6.5)),
    ("Week 8", "CFO, VP Procurement\nadopting independently", Inches(9.0)),
    ("Week 12", "HR, Product Dev\nusing AI organically", Inches(11.2)),
]

for label, desc, left in milestones:
    # Node circle
    add_shape(slide, left, timeline_y + Inches(0.15), Inches(0.16), Inches(0.16), ACCENT)
    # add a larger invisible shape behind to make node visible
    dot = add_shape(slide, left - Inches(0.04), timeline_y + Inches(0.11), Inches(0.24), Inches(0.24), ACCENT)
    # Label above
    add_text_box(slide, left - Inches(0.5), timeline_y - Inches(0.4), Inches(1.2), Inches(0.4),
                 label, 13, ACCENT, True, PP_ALIGN.CENTER)
    # Description below
    add_text_box(slide, left - Inches(0.6), timeline_y + Inches(0.6), Inches(1.5), Inches(0.8),
                 desc, 12, LIGHT_GRAY, False, PP_ALIGN.CENTER)

# Bottom section: what we did
add_shape(slide, Inches(0.8), Inches(4.5), Inches(11.7), Inches(2.5), RGBColor(0x3A, 0x3A, 0x3A))

tf = add_text_box(slide, Inches(1.3), Inches(4.7), Inches(5.0), Inches(2.0),
             "What we did:", 18, WHITE, True)
add_para(tf, "", 6, WHITE)
add_para(tf, "• KPI dashboard from raw financials", 14, LIGHT_GRAY)
add_para(tf, "• AI deployed across C-suite — COO, CFO, VP Procurement", 14, LIGHT_GRAY)
add_para(tf, "• Automated P&L analysis with exec summaries", 14, LIGHT_GRAY)
add_para(tf, "• Adoption spread organically to HR, Product Dev, Finance", 14, LIGHT_GRAY)

tf2 = add_text_box(slide, Inches(7.0), Inches(4.7), Inches(5.0), Inches(2.0),
             "< 3 months", 32, WHITE, True)
add_para(tf2, "Skeptical to self-sufficient", 14, MUTED)

# ═══════════════════════════════════════════
# SLIDE 8: PULL QUOTE (full slide)
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK)
add_bg_texture(slide)

add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10.3), Inches(0.8),
             "\u201C", 120, ACCENT, False, PP_ALIGN.LEFT)

add_text_box(slide, Inches(2.0), Inches(2.5), Inches(9.3), Inches(2.0),
             "I am really excited about what\nI just whipped together.", 44, WHITE, True, PP_ALIGN.CENTER)

add_shape(slide, Inches(6.0), Inches(5.0), Inches(1.3), Inches(0.05), ACCENT)

add_text_box(slide, Inches(2.0), Inches(5.4), Inches(9.3), Inches(0.6),
             "COO, Randy's Worldwide", 18, MUTED, False, PP_ALIGN.CENTER)
add_text_box(slide, Inches(2.0), Inches(5.9), Inches(9.3), Inches(0.5),
             "After 8 weeks of working with Day Horrigan", 14, MUTED, False, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# SLIDE 9: WHO WE ARE
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK)
add_bg_texture(slide)

add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(1.4), CREAM)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
             "Who we are", 32, DARK, True)

# Logan
add_shape(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5), RGBColor(0x3A, 0x3A, 0x3A))
add_shape(slide, Inches(0.8), Inches(2.0), Inches(0.08), Inches(4.5), ACCENT)
tf = add_text_box(slide, Inches(1.3), Inches(2.3), Inches(4.5), Inches(4.0),
             "Logan Day", 24, WHITE, True)
add_para(tf, "Darden MBA • U.S. Army Captain", 14, LIGHT_ACCENT)
add_para(tf, "", 10, WHITE)
add_para(tf, "Finance at Walmart by day, where I build AI tools and automation for executive teams. Built the dashboards and AI workflows at Randy's Worldwide.", 15, LIGHT_GRAY)
add_para(tf, "", 10, WHITE)
add_para(tf, "This is a focused side practice — we do this because we're genuinely good at it and the demand found us.", 15, LIGHT_GRAY)

# Matt
add_shape(slide, Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.5), RGBColor(0x3A, 0x3A, 0x3A))
add_shape(slide, Inches(7.0), Inches(2.0), Inches(0.08), Inches(4.5), ACCENT)
tf = add_text_box(slide, Inches(7.5), Inches(2.3), Inches(4.5), Inches(4.0),
             "Matt Horrigan", 24, WHITE, True)
add_para(tf, "Darden MBA • Strategy & Finance", 14, LIGHT_ACCENT)
add_para(tf, "", 10, WHITE)
add_para(tf, "Strategy and finance at WEX, with corporate strategy experience at Comcast.", 15, LIGHT_GRAY)
add_para(tf, "", 10, WHITE)
add_para(tf, "Brings the strategic lens — he sees the portco P&L the way your partners do.", 15, LIGHT_GRAY)

# ═══════════════════════════════════════════
# SLIDE 10: CTA
# ═══════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, CREAM)

# Natural Capital logo (subtle, bottom corner)
slide.shapes.add_picture(
    f"{ASSETS}/nc_logo.png",
    Inches(10.5), Inches(6.0),
    Inches(2.0), Inches(0.87)
)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(1.2),
             "30 minutes.", 52, DARK, True, PP_ALIGN.CENTER)

add_shape(slide, Inches(6.0), Inches(3.0), Inches(1.3), Inches(0.05), ACCENT)

add_text_box(slide, Inches(2.0), Inches(3.5), Inches(9.3), Inches(0.6),
             "That's all it takes to find the quick wins.", 20, MUTED, False, PP_ALIGN.CENTER)

tf = add_text_box(slide, Inches(2.0), Inches(4.5), Inches(9.3), Inches(2.5),
             "logan.day2018@gmail.com", 22, DARK, True, PP_ALIGN.CENTER)
add_para(tf, "", 14, DARK)
add_para(tf, "lj-consulting.vercel.app", 16, ACCENT, False, Pt(12), PP_ALIGN.CENTER)
add_para(tf, "", 10, DARK)
add_para(tf, "Logan Day  •  Matt Horrigan", 16, MUTED, False, Pt(12), PP_ALIGN.CENTER)

out = "/Users/loganday/Desktop/LJ_Consulting_x_Natural_Capital.pptx"
prs.save(out)
print(f"Saved to {out}")
